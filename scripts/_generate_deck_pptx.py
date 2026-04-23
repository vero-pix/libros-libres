"""Genera deck_valorizacion.pptx con los 11 slides del deck HTML."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Colores de marca
NAVY = RGBColor(0x1a, 0x1a, 0x2e)
NAVY_LIGHT = RGBColor(0x2a, 0x2a, 0x44)
GOLD = RGBColor(0xd4, 0xa0, 0x17)
GOLD_SOFT = RGBColor(0xe8, 0xc2, 0x57)
CREAM = RGBColor(0xfa, 0xf7, 0xf1)
INK = RGBColor(0x15, 0x15, 0x22)
MUTED = RGBColor(0x8b, 0x8b, 0x9c)
WHITE = RGBColor(0xff, 0xff, 0xff)

# Fuentes
FONT_SERIF = 'Georgia'
FONT_SANS = 'Helvetica Neue'

# Tamaño 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def blank_slide(bg_color):
    """Crea un slide en blanco con fondo sólido."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height,
             font=FONT_SANS, size=16, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    # Si el texto tiene \n, split en runs
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            para = p
        else:
            para = tf.add_paragraph()
            para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_tag(slide, tag_text, color=GOLD):
    """Tag pequeño arriba izquierda."""
    add_text(slide, tag_text.upper(), Inches(0.6), Inches(0.5),
             Inches(6), Inches(0.3),
             font=FONT_SANS, size=10, bold=True, color=color)


def add_page_num(slide, num, color=MUTED):
    add_text(slide, num, Inches(12), Inches(7.0),
             Inches(1.1), Inches(0.3),
             font=FONT_SANS, size=9, color=color, align=PP_ALIGN.RIGHT)


def add_divider(slide, left, top, width=Inches(0.7), color=GOLD):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_pill(slide, text, left, top, bg=GOLD, text_color=NAVY):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   Inches(4.2), Inches(0.35))
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg
    pill.line.fill.background()
    pill.adjustments[0] = 0.5
    tf = pill.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_SANS
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_color


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ═══════════════════════════════════════════════════════════════
s = blank_slide(NAVY)
add_pill(s, "Propuesta de Inversión · Abril 2026", Inches(0.6), Inches(0.5))

# Marca "tl" circle
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.2),
                              Inches(0.6), Inches(0.6))
circle.fill.background()
circle.line.color.rgb = GOLD
circle.line.width = Pt(2)
ctf = circle.text_frame
ctf.margin_left = 0; ctf.margin_right = 0
ctf.margin_top = 0; ctf.margin_bottom = 0
cp = ctf.paragraphs[0]
cp.alignment = PP_ALIGN.CENTER
crun = cp.add_run()
crun.text = "tl"
crun.font.name = FONT_SERIF
crun.font.size = Pt(22)
crun.font.color.rgb = GOLD

add_text(s, "tuslibros.cl", Inches(0.6), Inches(2.3),
         Inches(12), Inches(2.3),
         font=FONT_SERIF, size=92, bold=False, color=CREAM)

add_text(s, "Los libros de tu ciudad.\nMarketplace de libros usados con mapa, curación y split payment.",
         Inches(0.6), Inches(4.7), Inches(11), Inches(1.5),
         font=FONT_SANS, size=20, color=CREAM)

add_text(s, "VERÓNICA VELÁSQUEZ · FOUNDER", Inches(0.6), Inches(6.9),
         Inches(6), Inches(0.3),
         font=FONT_SANS, size=10, bold=True, color=MUTED)
add_text(s, "SANTIAGO, CHILE · ABRIL 2026", Inches(7), Inches(6.9),
         Inches(5.8), Inches(0.3),
         font=FONT_SANS, size=10, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
add_page_num(s, "01 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problema
# ═══════════════════════════════════════════════════════════════
s = blank_slide(CREAM)
add_tag(s, "El problema")
add_text(s, "El mercado de libros usados en Chile es grande,\npero está invisible.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(1.8),
         font=FONT_SERIF, size=42, color=INK)
add_divider(s, Inches(0.6), Inches(2.9))

# 3 métricas en grid
metrics = [
    ("$8-12", "mil millones CLP / año", "Mercado estimado de libros usados en Chile, repartido entre MercadoLibre, Facebook Marketplace, librerías de viejo e intercambios informales. ~USD 8-12M."),
    ("0", "Plataforma digital dedicada", "Buscalibre domina libro nuevo. MercadoLibre es commodity sin curación. Facebook no tiene pago integrado ni confianza. Las librerías de viejo siguen offline."),
    ("72%", "Catálogo invisible", "Antes de tuslibros: el libro usado que quieres existe, pero encontrarlo toma tres búsquedas, cero garantía y una conversación de WhatsApp con un desconocido."),
]
col_w = Inches(3.9)
for i, (num, lbl, desc) in enumerate(metrics):
    left = Inches(0.6 + i * 4.15)
    add_text(s, num, left, Inches(3.4), col_w, Inches(1.4),
             font=FONT_SERIF, size=72, color=GOLD)
    add_text(s, lbl.upper(), left, Inches(4.9), col_w, Inches(0.4),
             font=FONT_SANS, size=10, bold=True, color=MUTED)
    add_text(s, desc, left, Inches(5.4), col_w, Inches(1.8),
             font=FONT_SANS, size=11, color=INK)

add_text(s, "En Chile compramos ~45 millones de libros al año. El mercado de segunda vida existe, pero no tiene infraestructura digital. Nadie lo está resolviendo.",
         Inches(0.6), Inches(6.85), Inches(11), Inches(0.5),
         font=FONT_SANS, size=11, italic=True, color=MUTED)
add_page_num(s, "02 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Solución
# ═══════════════════════════════════════════════════════════════
s = blank_slide(NAVY)
add_tag(s, "La solución")
add_text(s, "Un marketplace con tres diferenciadores\nque nadie tiene junto.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(1.8),
         font=FONT_SERIF, size=38, color=GOLD_SOFT)

diffs = [
    ("Mapa geolocalizado", "Ves los libros como puntos en el mapa de tu ciudad. El libro más cercano es el mejor resultado. Experiencia tipo Uber, no catálogo plano."),
    ("Curación de rarezas", "Featured, coleccionables y top 10 priorizan primeras ediciones, títulos agotados, editoriales históricas. No competimos en commodity — competimos en lo que nadie tiene."),
    ("Split payment integrado", "OAuth con MercadoPago: el pago se divide automáticamente entre vendedor y plataforma. Courier Shipit genera etiqueta al confirmar el pago. Un solo flujo, cero fricción operativa."),
]
for i, (title, desc) in enumerate(diffs):
    left = Inches(0.6 + i * 4.15)
    add_text(s, title.upper(), left, Inches(3.2), Inches(3.9), Inches(0.4),
             font=FONT_SANS, size=12, bold=True, color=GOLD)
    add_text(s, desc, left, Inches(3.75), Inches(3.9), Inches(2.5),
             font=FONT_SANS, size=13, color=CREAM)

# Caja destacada "foto-a-catálogo"
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.9),
                          Inches(12.1), Inches(1.1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x2a, 0x22, 0x18)
box.line.fill.background()
# Borde izquierdo dorado
border = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.9),
                             Inches(0.04), Inches(1.1))
border.fill.solid()
border.fill.fore_color.rgb = GOLD
border.line.fill.background()

add_text(s, "FEATURE ÚNICO EN CHILE: FOTO-A-CATÁLOGO",
         Inches(0.9), Inches(6.02), Inches(11), Inches(0.3),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
add_text(s, "Subes la foto de la portada. Identificación vía LLM + enriquecimiento automático con Google Books y Open Library. Ningún competidor en Chile tiene esto. Elimina el 90% del trabajo de publicar.",
         Inches(0.9), Inches(6.4), Inches(11.8), Inches(0.55),
         font=FONT_SANS, size=12, color=CREAM)
add_page_num(s, "03 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Producto
# ═══════════════════════════════════════════════════════════════
s = blank_slide(CREAM)
add_tag(s, "Producto")
add_text(s, "En producción real desde el Viernes Santo.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
         font=FONT_SERIF, size=36, color=INK)
add_text(s, "Construido en 7 días. Iterado 18.",
         Inches(0.6), Inches(1.85), Inches(12), Inches(0.7),
         font=FONT_SERIF, size=28, italic=True, color=MUTED)
add_divider(s, Inches(0.6), Inches(2.7))

# 2 columnas
stack_items = [
    "Next.js 14 App Router + React 18 + TypeScript",
    "Supabase (Postgres + Storage + Auth)",
    "MercadoPago OAuth con split payment",
    "Shipit (Starken, Chilexpress, Blue Express, 99 Minutos)",
    "Mapbox GL + geolocalización precisa",
    "Vercel + CI/CD automático en producción",
]
func_items = [
    "Bundle checkout multi-libro con preferencia MP única",
    "URLs amigables /libro/[vendedor]/[slug]",
    "Mensajería interna con RLS",
    "Sistema de categorías con árbol (5 raíces · 32 subcategorías)",
    "Novedades, destacados, coleccionables, tiers",
    "Admin panel con CRUD completo",
]

add_text(s, "STACK MODERNO", Inches(0.6), Inches(3.0),
         Inches(5.5), Inches(0.4),
         font=FONT_SANS, size=12, bold=True, color=GOLD)
for i, item in enumerate(stack_items):
    add_text(s, "→  " + item, Inches(0.6), Inches(3.5 + i * 0.35),
             Inches(6), Inches(0.35),
             font=FONT_SANS, size=12, color=INK)

add_text(s, "FUNCIONALIDAD OPERATIVA", Inches(6.9), Inches(3.0),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=12, bold=True, color=GOLD)
for i, item in enumerate(func_items):
    add_text(s, "→  " + item, Inches(6.9), Inches(3.5 + i * 0.35),
             Inches(6.2), Inches(0.35),
             font=FONT_SANS, size=12, color=INK)

add_text(s, "Costo de reproducir esto con una agencia de desarrollo: USD 150k-250k. Costo real construido con AI-augmented development: el equivalente a 7 días-persona.",
         Inches(0.6), Inches(6.85), Inches(12), Inches(0.5),
         font=FONT_SANS, size=10, italic=True, color=MUTED)
add_page_num(s, "04 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Tracción
# ═══════════════════════════════════════════════════════════════
s = blank_slide(NAVY)
add_tag(s, "Tracción · 18 días")
add_text(s, "Sin distribución activa ni gasto en marketing.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
         font=FONT_SERIF, size=34, color=GOLD_SOFT)

# 8 métricas en 2 filas de 4
metrics_traccion = [
    ("438", "Usuarios activos", "GA4, primeros 10 días"),
    ("3", "Compradores reales", "11 libros vendidos"),
    ("3", "Librerías activadas", "Orgánicamente, sin outreach"),
    ("207", "Listings curados", "Valor catálogo $2.4M CLP"),
    ("20%", "Conversión registro → compra", "Benchmark seed: 5-15%"),
    ("4.6", "Vistas / usuario", "Engagement sano"),
    ("$90k", "GMV acumulado CLP", "3 ventas en 18 días"),
    ("$45", "Burn mensual USD", "Vercel Pro + Supabase Pro"),
]
for i, (num, lbl, sub) in enumerate(metrics_traccion):
    row = i // 4
    col = i % 4
    left = Inches(0.6 + col * 3.1)
    top = Inches(2.5 + row * 2.1)
    add_text(s, num, left, top, Inches(3), Inches(1.1),
             font=FONT_SERIF, size=54, color=GOLD)
    add_text(s, lbl.upper(), left, top + Inches(1.1), Inches(3), Inches(0.4),
             font=FONT_SANS, size=10, bold=True, color=RGBColor(0xcc, 0xcc, 0xd5))
    add_text(s, sub, left, top + Inches(1.45), Inches(3), Inches(0.4),
             font=FONT_SANS, size=10, italic=True, color=RGBColor(0xa0, 0xa0, 0xb0))

add_text(s, "Operación iniciada 3 abril 2026. Tres ventas reales en 18 días: Zdravko (La Marina Tomo I, $5.000), set de 7 Maigret a tercer comprador ($40.000), bundle multi-libro a Camilo (17 abril, Concepción, $45.000).",
         Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         font=FONT_SANS, size=9, italic=True, color=RGBColor(0x8b, 0x8b, 0x9c))
add_page_num(s, "05 / 11", color=RGBColor(0x8b, 0x8b, 0x9c))


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Evidencia / Testimonio
# ═══════════════════════════════════════════════════════════════
s = blank_slide(CREAM)
add_tag(s, "Evidencia")
add_text(s, "El patrón ya está validado.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
         font=FONT_SERIF, size=38, color=INK)

# Quote
quote_left = Inches(0.6)
quote_top = Inches(2.2)
# Línea vertical dorada
qline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, quote_left, quote_top,
                            Inches(0.05), Inches(1.8))
qline.fill.solid()
qline.fill.fore_color.rgb = GOLD
qline.line.fill.background()

add_text(s, '"Fácil y sin complicaciones. Muy buena disposición por parte del vendedor,\nvolvería a comprar sin ningún problema."',
         Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.6),
         font=FONT_SERIF, size=30, italic=True, color=INK)
add_text(s, "— Z., PRIMER COMPRADOR · LA MARINA EN LA HISTORIA DE CHILE · 7 ABRIL 2026",
         Inches(1.0), Inches(3.85), Inches(11.5), Inches(0.3),
         font=FONT_SANS, size=10, bold=True, color=MUTED)

# 2 columnas abajo
add_text(s, "QUÉ VENDEMOS", Inches(0.6), Inches(4.7), Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
add_text(s, "Libros raros a precio 40-60% bajo IberLibro/Buscalibre. El bundle de Camilo incluyó Parra ($23.992, IberLibro €30-50) y De Rokha edición cubana 1991 agotada ($14.000). Son joyas, no commodity.",
         Inches(0.6), Inches(5.15), Inches(6), Inches(1.7),
         font=FONT_SANS, size=12, color=INK)

add_text(s, "POR QUÉ NO MORIREMOS EN COMMODITY", Inches(6.9), Inches(4.7),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
add_text(s, "Quien llega a tuslibros ya filtró Buscalibre y MercadoLibre. Busca lo que esos no tienen. Competir en best-seller moderno es suicidio. La ventaja defendible es la curación.",
         Inches(6.9), Inches(5.15), Inches(6), Inches(1.7),
         font=FONT_SANS, size=12, color=INK)

add_text(s, "Publicados 3 libros el 16 abril a las 21h. Vendidos en bundle 14 horas después, 17 abril 15h. Tiempo entre publicación y venta en caso de libro raro + precio competitivo: menos de un día.",
         Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         font=FONT_SANS, size=10, italic=True, color=MUTED)
add_page_num(s, "06 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Mercado
# ═══════════════════════════════════════════════════════════════
s = blank_slide(NAVY)
add_tag(s, "Mercado")
add_text(s, "TAM, SAM, SOM — conservador.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
         font=FONT_SERIF, size=38, color=GOLD_SOFT)

# Caja de definición TAM/SAM/SOM
definition_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0),
                                      Inches(12.1), Inches(0.6))
definition_box.fill.solid()
definition_box.fill.fore_color.rgb = RGBColor(0x2a, 0x22, 0x18)
definition_box.line.fill.background()
dborder = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0),
                              Inches(0.04), Inches(0.6))
dborder.fill.solid()
dborder.fill.fore_color.rgb = GOLD
dborder.line.fill.background()
add_text(s, "TAM (Total Addressable): techo teórico  ·  SAM (Serviceable Addressable): lo que el producto realmente puede servir  ·  SOM (Serviceable Obtainable): lo realista capturar en 3 años",
         Inches(0.9), Inches(2.15), Inches(11.8), Inches(0.35),
         font=FONT_SANS, size=11, color=CREAM)

market = [
    ("$10 mil millones", "CLP / año · TAM Chile", "Total de libros usados transados en Chile al año entre todos los canales. ~USD 10.5M."),
    ("$3-4 mil millones", "CLP / año · SAM urbano", "Santiago, Valparaíso, Concepción. Donde el mapa y courier funcionan. ~USD 3.2-4.2M."),
    ("$300-400 millones", "CLP / año · SOM año 3", "Con 10% de penetración digital del mercado urbano, GMV objetivo. ~USD 315-420k."),
]
for i, (num, lbl, desc) in enumerate(market):
    left = Inches(0.6 + i * 4.15)
    add_text(s, num, left, Inches(2.9), Inches(3.9), Inches(1.2),
             font=FONT_SERIF, size=38, color=GOLD)
    add_text(s, lbl.upper(), left, Inches(4.15), Inches(3.9), Inches(0.5),
             font=FONT_SANS, size=11, bold=True, color=RGBColor(0xcc, 0xcc, 0xd5))
    add_text(s, desc, left, Inches(4.7), Inches(3.9), Inches(1.3),
             font=FONT_SANS, size=11, color=CREAM)

add_text(s, "CRECIMIENTO ESTRUCTURAL", Inches(0.6), Inches(5.85),
         Inches(12), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)

struct = [
    "Mercado del libro físico en Chile crece 8-12% anual post-pandemia (Cámara Chilena del Libro)",
    "Economía circular + sustentabilidad empuja usado vs. nuevo",
    "Generación universitaria busca libros de cátedra usados (nicho de alta recurrencia)",
    "Coleccionismo en alza: bibliófilos chilenos pagan CLP 50-300k por primera edición nacional",
]
for i, item in enumerate(struct):
    add_text(s, "→  " + item, Inches(0.6), Inches(6.3 + i * 0.27),
             Inches(12), Inches(0.25),
             font=FONT_SANS, size=10, color=CREAM)
add_page_num(s, "07 / 11", color=RGBColor(0x8b, 0x8b, 0x9c))


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Competencia
# ═══════════════════════════════════════════════════════════════
s = blank_slide(CREAM)
add_tag(s, "Competencia")
add_text(s, "No competimos contra ellos —\nresolvemos lo que ninguno resuelve.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(1.6),
         font=FONT_SERIF, size=32, color=INK)

# Tabla
headers = ["Jugador", "Libro nuevo", "Libro usado", "Curación", "Mapa", "Split pay."]
rows = [
    ["Buscalibre", "✓", "—", "Editorial", "—", "—"],
    ["MercadoLibre", "✓", "commodity", "—", "—", "Parcial"],
    ["Facebook MP", "—", "informal", "—", "—", "—"],
    ["Librerías viejo", "—", "offline", "Curador", "—", "—"],
    ["tuslibros.cl", "—", "curado", "featured", "✓", "✓ OAuth MP"],
]

table_left = Inches(0.6)
table_top = Inches(3.1)
table_w = Inches(12.1)
col_widths = [Inches(2.5), Inches(1.5), Inches(1.9), Inches(2.0), Inches(1.4), Inches(2.8)]

# Headers
x = table_left
for j, h in enumerate(headers):
    add_text(s, h.upper(), x, table_top, col_widths[j], Inches(0.35),
             font=FONT_SANS, size=10, bold=True, color=GOLD)
    x += col_widths[j]

# Line under headers
add_divider(s, table_left, table_top + Inches(0.45), width=table_w, color=MUTED)

# Rows
for i, row in enumerate(rows):
    row_top = table_top + Inches(0.6 + i * 0.55)
    is_last = (i == len(rows) - 1)
    # Highlight last row
    if is_last:
        hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, table_left, row_top - Inches(0.07),
                                 table_w, Inches(0.5))
        hl.fill.solid()
        hl.fill.fore_color.rgb = RGBColor(0xf0, 0xe4, 0xc8)
        hl.line.fill.background()

    x = table_left
    for j, cell in enumerate(row):
        bold = is_last or j == 0
        color = INK
        add_text(s, cell, x, row_top, col_widths[j], Inches(0.4),
                 font=FONT_SANS, size=12, bold=bold, color=color)
        x += col_widths[j]

add_text(s, "La posición competitiva no es \"otro Buscalibre\" ni \"otro MercadoLibre\". Es la librería de viejo de Santiago, pero online, con mapa, con pago integrado y con foto-a-catálogo.",
         Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
         font=FONT_SANS, size=11, italic=True, color=MUTED)
add_page_num(s, "08 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Equipo
# ═══════════════════════════════════════════════════════════════
s = blank_slide(NAVY)
add_tag(s, "Equipo")
add_text(s, "Founder única, velocity probada.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
         font=FONT_SERIF, size=38, color=GOLD_SOFT)

add_text(s, "VERÓNICA VELÁSQUEZ — FOUNDER", Inches(0.6), Inches(2.6),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
add_text(s, "Background en economía y consultoría. Ex-economics.cl. Decisión deliberada de operar sin cofundador técnico.",
         Inches(0.6), Inches(3.1), Inches(5.8), Inches(1.3),
         font=FONT_SANS, size=13, color=CREAM)
add_text(s, "No es limitación: es el modelo. Dev asistido por Claude Code me permite iterar producción en horas lo que una startup tradicional iteraría en semanas.",
         Inches(0.6), Inches(4.5), Inches(5.8), Inches(1.8),
         font=FONT_SANS, size=13, italic=True, color=CREAM)

add_text(s, "EVIDENCIA DE EJECUCIÓN", Inches(6.9), Inches(2.6),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
evidence = [
    "7 días desde 0 hasta producción con MP + Shipit",
    "18 días de iteración continua: 207 listings curados, 3 librerías onboardeadas, 3 compradores, 11 libros vendidos",
    "Deploy diario a producción sin regresiones críticas",
    "Burn $45/mes: disciplina de capital desde día uno",
    "Founder cerca del customer: respondí yo el primer email, yo cerré el primer testimonio, yo entregué el primer bundle",
]
for i, item in enumerate(evidence):
    add_text(s, "→  " + item, Inches(6.9), Inches(3.1 + i * 0.65),
             Inches(6.2), Inches(0.6),
             font=FONT_SANS, size=11, color=CREAM)

add_text(s, "El argumento tradicional contra un founder único es ejecución. La evidencia de los últimos 18 días responde eso antes de la pregunta.",
         Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         font=FONT_SANS, size=10, italic=True, color=RGBColor(0xa0, 0xa0, 0xb0))
add_page_num(s, "09 / 11", color=RGBColor(0x8b, 0x8b, 0x9c))


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Propuesta / Valorización
# ═══════════════════════════════════════════════════════════════
s = blank_slide(CREAM)
add_tag(s, "Propuesta")
add_text(s, "USD 80-120k · 6-9% · pre-money USD 1.2M.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
         font=FONT_SERIF, size=34, color=INK)
add_text(s, "$76 – $114 millones CLP  ·  ronda pre-seed  ·  SAFE o nota convertible  ·  runway 12 meses",
         Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
         font=FONT_SANS, size=13, italic=True, color=MUTED)
add_divider(s, Inches(0.6), Inches(2.15))

# Columna izquierda — valorización
add_text(s, "VALORIZACIÓN TRIANGULADA", Inches(0.6), Inches(2.45),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)

val_rows = [
    ("Berkus Method", "USD 1.6M"),
    ("Scorecard vs. comparables LatAm", "USD 1.42M"),
    ("Replacement cost (piso)", "USD 200k"),
    ("Comparables Chile seed marketplace", "USD 1.2-1.4M"),
    ("Pre-money objetivo", "USD 1.2M"),
]
for i, (lbl, val) in enumerate(val_rows):
    top = Inches(2.9 + i * 0.38)
    bold = (i == len(val_rows) - 1)
    if bold:
        hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top - Inches(0.04),
                                 Inches(6), Inches(0.4))
        hl.fill.solid()
        hl.fill.fore_color.rgb = RGBColor(0xf0, 0xe4, 0xc8)
        hl.line.fill.background()
    add_text(s, lbl, Inches(0.6), top, Inches(4), Inches(0.35),
             font=FONT_SANS, size=11, bold=bold, color=INK)
    add_text(s, val, Inches(4.6), top, Inches(2), Inches(0.35),
             font=FONT_SANS, size=11, bold=bold, color=INK, align=PP_ALIGN.RIGHT)

# Caja métodos breve
methods_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.95),
                                   Inches(6), Inches(1.9))
methods_box.fill.solid()
methods_box.fill.fore_color.rgb = RGBColor(0xee, 0xe8, 0xdc)
methods_box.line.fill.background()
mborder = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.95),
                              Inches(0.04), Inches(1.9))
mborder.fill.solid()
mborder.fill.fore_color.rgb = GOLD
mborder.line.fill.background()

add_text(s, "MÉTODOS, EN BREVE", Inches(0.85), Inches(5.05),
         Inches(5.5), Inches(0.3),
         font=FONT_SANS, size=10, bold=True, color=GOLD)
methods_text = [
    ("Berkus", "asigna hasta USD 500k a 5 componentes (idea, producto, equipo, alianzas, ventas)."),
    ("Scorecard", "compara contra comparables seed LatAm, ajusta por fuerzas relativas."),
    ("Replacement cost", "cuánto costaría reconstruir lo construido. Piso de negociación."),
]
for i, (name, desc) in enumerate(methods_text):
    top = Inches(5.4 + i * 0.48)
    add_text(s, name, Inches(0.85), top, Inches(1.5), Inches(0.3),
             font=FONT_SANS, size=10, bold=True, color=INK)
    add_text(s, desc, Inches(2.35), top, Inches(4.2), Inches(0.45),
             font=FONT_SANS, size=10, color=INK)

# Columna derecha — uso de fondos
add_text(s, "USO DE FONDOS · 12 MESES  (USD 100k base)", Inches(6.9), Inches(2.45),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)

funds_rows = [
    ("Salario founder ($5M CLP/mes × 12)", "60%"),
    ("Marketing digital + GTM", "20%"),
    ("Activación librerías + features", "12%"),
    ("Infra + reserva operativa", "8%"),
]
for i, (lbl, val) in enumerate(funds_rows):
    top = Inches(2.9 + i * 0.4)
    add_text(s, lbl, Inches(6.9), top, Inches(5), Inches(0.35),
             font=FONT_SANS, size=11, color=INK,
             bold=(i == 0))
    add_text(s, val, Inches(11.9), top, Inches(0.8), Inches(0.35),
             font=FONT_SANS, size=11, bold=True, color=INK, align=PP_ALIGN.RIGHT)

add_text(s, "Dedicación full-time. Hoy opero con ahorros — el salario en el ask libera 100% del tiempo para escalar.",
         Inches(6.9), Inches(4.6), Inches(6), Inches(0.7),
         font=FONT_SANS, size=10, italic=True, color=MUTED)

add_text(s, "HITO A 12 MESES", Inches(6.9), Inches(5.4),
         Inches(6), Inches(0.4),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
hitos = [
    "20 vendedores activos · 2 alianzas universitarias",
    "$5 millones CLP MRR · 500+ transacciones/mes",
    "Listo para ronda seed USD 500k-1M",
]
for i, h in enumerate(hitos):
    add_text(s, "→  " + h, Inches(6.9), Inches(5.85 + i * 0.35),
             Inches(6), Inches(0.3),
             font=FONT_SANS, size=11, color=INK)

add_page_num(s, "10 / 11")


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Cierre
# ═══════════════════════════════════════════════════════════════
s = blank_slide(NAVY)
add_tag(s, "Hablemos")
add_text(s, "El mejor momento para",
         Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
         font=FONT_SERIF, size=60, color=CREAM)
add_text(s, "entrar es ahora.",
         Inches(0.6), Inches(2.4), Inches(12), Inches(1.2),
         font=FONT_SERIF, size=60, italic=True, color=GOLD)

add_text(s, "En 18 días el producto está en producción, hay compradores reales con ciclo cerrado, 3 librerías profesionales llegaron solas y el primer testimonio autorizado está publicado. El capital que busco es para activar lo que ya está tocando la puerta — no para construir algo desde cero.",
         Inches(0.6), Inches(4.2), Inches(12), Inches(1.5),
         font=FONT_SANS, size=15, color=CREAM)

# CTA box
cta = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.9),
                          Inches(12.1), Inches(1.3))
cta.fill.solid()
cta.fill.fore_color.rgb = RGBColor(0x22, 0x1e, 0x10)
cta.line.color.rgb = GOLD
cta.line.width = Pt(1.5)

add_text(s, "VERÓNICA VELÁSQUEZ — FOUNDER · TUSLIBROS.CL",
         Inches(0.9), Inches(6.05), Inches(8), Inches(0.35),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
add_text(s, "vero@economics.cl · Santiago · Providencia",
         Inches(0.9), Inches(6.45), Inches(8), Inches(0.35),
         font=FONT_SANS, size=12, color=CREAM)
add_text(s, "PRÓXIMOS PASOS",
         Inches(8.8), Inches(6.05), Inches(4), Inches(0.35),
         font=FONT_SANS, size=11, bold=True, color=GOLD)
add_text(s, "Deep dive GA4 · Proyección 36 meses · Term sheet 1-2 semanas",
         Inches(8.8), Inches(6.45), Inches(4.5), Inches(0.7),
         font=FONT_SANS, size=11, color=CREAM)

add_page_num(s, "11 / 11", color=RGBColor(0x8b, 0x8b, 0x9c))


# Guardar
output = '/Users/veronicavelasquez/libros-libres/docs/deck_valorizacion.pptx'
prs.save(output)
print(f"✓ Generado: {output}")
print(f"  {len(prs.slides)} slides, 16:9 widescreen")
